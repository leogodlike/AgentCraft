#!/usr/bin/env python3
"""Generate a Huawei-style PPTX about Machine Learning using python-pptx"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# Huawei brand colors
HUAWEI_RED = RGBColor(0xC7, 0x02, 0x0E)
HUAWEI_RED_LIGHT = RGBColor(0xD5, 0x3C, 0x44)
HUAWEI_RED_BG = RGBColor(0xFF, 0xF1, 0xEF)
HUAWEI_GRAY1 = RGBColor(0x8C, 0x8C, 0x8C)
HUAWEI_GRAY3 = RGBColor(0xBF, 0xBF, 0xBF)
HUAWEI_GRAY4 = RGBColor(0xD9, 0xD9, 0xD9)
HUAWEI_GRAY5 = RGBColor(0xF2, 0xF2, 0xF2)
BLACK = RGBColor(0x00, 0x00, 0x00)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

def add_white_bg(slide):
    """Set white background"""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = WHITE

def add_red_top_line(slide):
    """Add 1px red top line"""
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Pt(1.5)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = HUAWEI_RED
    line.line.fill.background()

def add_title_bar(slide, cn_title, en_title, page_num):
    """Standard title bar with red indicator"""
    # Title bar background
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, Pt(1.5), SLIDE_W, Inches(1.25)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = WHITE
    bar.line.color.rgb = HUAWEI_GRAY4
    bar.line.width = Pt(1)
    
    # Red indicator
    ind = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(0.45), Pt(4), Inches(0.42)
    )
    ind.fill.solid()
    ind.fill.fore_color.rgb = HUAWEI_RED
    ind.line.fill.background()
    
    # Title
    txBox = slide.shapes.add_textbox(Inches(0.75), Inches(0.35), Inches(8), Inches(0.6))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = cn_title
    p.font.size = Pt(35)
    p.font.bold = True
    p.font.color.rgb = BLACK
    
    # English subtitle
    txBox2 = slide.shapes.add_textbox(Inches(0.75), Inches(0.78), Inches(6), Inches(0.4))
    tf2 = txBox2.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = en_title
    p2.font.size = Pt(19)
    p2.font.color.rgb = HUAWEI_GRAY1
    
    # Page number
    txBox3 = slide.shapes.add_textbox(Inches(11.5), Inches(0.45), Inches(1.5), Inches(0.5))
    tf3 = txBox3.text_frame
    p3 = tf3.paragraphs[0]
    p3.text = page_num
    p3.font.size = Pt(19)
    p3.font.color.rgb = HUAWEI_GRAY1
    p3.alignment = PP_ALIGN.RIGHT

# ============ PAGE 1: COVER ============
slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # blank
bg = slide1.background
fill = bg.fill
fill.solid()
fill.fore_color.rgb = HUAWEI_RED

# Decorative text
txBox = slide1.shapes.add_textbox(Inches(9), Inches(1.5), Inches(4), Inches(4))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "AI"
p.font.size = Pt(300)
p.font.bold = True
p.font.color.rgb = WHITE
p.font.color.brightness = 0.94

# English top label
txBox = slide1.shapes.add_textbox(Inches(1.2), Inches(3.0), Inches(6), Inches(0.5))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "ARTIFICIAL INTELLIGENCE"
p.font.size = Pt(18)
p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
p.font.color.brightness = 0.8

# Main title
txBox = slide1.shapes.add_textbox(Inches(1.2), Inches(3.6), Inches(10), Inches(1.5))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "什么是机器学习？"
p.font.size = Pt(62)
p.font.bold = True
p.font.color.rgb = WHITE

# Subtitle
txBox = slide1.shapes.add_textbox(Inches(1.2), Inches(5.0), Inches(8), Inches(0.6))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "Machine Learning Overview"
p.font.size = Pt(24)
p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
p.font.color.brightness = 0.7

# Decorative line
rect = slide1.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Inches(1.2), Inches(5.6), Inches(1.5), Pt(3)
)
rect.fill.solid()
rect.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
rect.fill.fore_color.brightness = 0.5
rect.line.fill.background()

# Bottom text
txBox = slide1.shapes.add_textbox(Inches(1.2), Inches(6.5), Inches(6), Inches(0.4))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "人工智能核心分支 · 数据驱动的智能革命"
p.font.size = Pt(16)
p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
p.font.color.brightness = 0.5

txBox = slide1.shapes.add_textbox(Inches(11.5), Inches(6.5), Inches(1.5), Inches(0.4))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "01"
p.font.size = Pt(16)
p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
p.font.color.brightness = 0.5
p.alignment = PP_ALIGN.RIGHT

# ============ PAGE 2: CONTENTS ============
slide2 = prs.slides.add_slide(prs.slide_layouts[6])
add_white_bg(slide2)
add_red_top_line(slide2)
add_title_bar(slide2, "目录", "CONTENTS", "02")

# Contents items - two columns
left_items = [
    ("01", "什么是机器学习"),
    ("02", "核心思想与对比"),
    ("03", "机器学习基本流程"),
    ("04", "三大学习类型"),
]
right_items = [
    ("05", "监督学习 & 无监督学习"),
    ("06", "强化学习"),
    ("07", "日常应用场景"),
    ("08", "总结与展望"),
]

def add_content_item(slide, num, title, left, top, width, height):
    """Add a content item with number"""
    # Background card
    card = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, height
    )
    card.fill.solid()
    card.fill.fore_color.rgb = HUAWEI_GRAY5
    card.line.fill.background()
    
    # Left red border
    border = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, Pt(4), height
    )
    border.fill.solid()
    border.fill.fore_color.rgb = HUAWEI_RED if num == "01" else HUAWEI_GRAY4
    border.line.fill.background()
    
    # Number
    txBox = slide.shapes.add_textbox(left + Inches(0.3), top + Inches(0.15), Inches(0.6), height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = num
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = HUAWEI_RED
    
    # Title
    txBox = slide.shapes.add_textbox(left + Inches(0.9), top + Inches(0.15), width - Inches(1.2), height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(21)
    p.font.bold = True
    p.font.color.rgb = BLACK
    p.alignment = PP_ALIGN.LEFT

item_w = Inches(5.2)
item_h = Inches(0.6)
start_left = Inches(0.8)
start_top = Inches(1.8)
gap = Inches(0.25)

for i, (num, title) in enumerate(left_items):
    add_content_item(slide2, num, title, start_left, start_top + i * (item_h + gap), item_w, item_h)

for i, (num, title) in enumerate(right_items):
    add_content_item(slide2, num, title, start_left + Inches(6.5), start_top + i * (item_h + gap), item_w, item_h)

# ============ PAGE 3: What is ML ============
slide3 = prs.slides.add_slide(prs.slide_layouts[6])
add_white_bg(slide3)
add_red_top_line(slide3)
add_title_bar(slide3, "什么是机器学习", "What is Machine Learning", "03")

# Definition card
card = slide3.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(1.6), Inches(12.1), Inches(1.6)
)
card.fill.solid()
card.fill.fore_color.rgb = HUAWEI_RED_BG
card.line.color.rgb = HUAWEI_RED
card.line.width = Pt(1)

# ML label in card
label = slide3.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.8), Inches(0.55), Inches(0.55)
)
label.fill.solid()
label.fill.fore_color.rgb = HUAWEI_RED
label.line.fill.background()
txBox = slide3.shapes.add_textbox(Inches(0.83), Inches(1.83), Inches(0.5), Inches(0.5))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "ML"
p.font.size = Pt(16)
p.font.bold = True
p.font.color.rgb = WHITE
p.alignment = PP_ALIGN.CENTER

# Definition title
txBox = slide3.shapes.add_textbox(Inches(1.5), Inches(1.7), Inches(4), Inches(0.4))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "核心定义"
p.font.size = Pt(21)
p.font.bold = True
p.font.color.rgb = HUAWEI_RED

# Definition text
txBox = slide3.shapes.add_textbox(Inches(1.5), Inches(2.2), Inches(10.5), Inches(0.9))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = '机器学习是人工智能（AI）的核心分支，它让计算机不需要被明确编程就能从数据中学习和改进。简单来说，传统编程是你告诉计算机\u201c怎么做\u201d，而机器学习是你给计算机数据和答案，让它自己找出\u201c怎么做\u201d的规律。'
p.font.size = Pt(17)
p.font.color.rgb = BLACK

# Comparison title
txBox = slide3.shapes.add_textbox(Inches(0.6), Inches(3.5), Inches(8), Inches(0.4))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "传统编程 vs 机器学习：两种截然不同的范式"
p.font.size = Pt(20)
p.font.bold = True
p.font.color.rgb = BLACK

# Traditional programming card
card_t = slide3.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(4.0), Inches(5.8), Inches(3.0)
)
card_t.fill.solid()
card_t.fill.fore_color.rgb = WHITE
card_t.line.color.rgb = HUAWEI_GRAY4
card_t.line.width = Pt(1)

header_t = slide3.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(4.0), Inches(5.8), Inches(0.5)
)
header_t.fill.solid()
header_t.fill.fore_color.rgb = HUAWEI_GRAY1
header_t.line.fill.background()
txBox = slide3.shapes.add_textbox(Inches(0.6), Inches(4.05), Inches(5.8), Inches(0.4))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "传统编程"
p.font.size = Pt(18)
p.font.bold = True
p.font.color.rgb = WHITE
p.alignment = PP_ALIGN.CENTER

# Content
txBox = slide3.shapes.add_textbox(Inches(1.0), Inches(4.7), Inches(5.0), Inches(2.0))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "输入：数据 + 规则"
p.font.size = Pt(20)
p.font.bold = True
p.font.color.rgb = HUAWEI_RED
p.alignment = PP_ALIGN.CENTER
p2 = tf.add_paragraph()
p2.text = "↓"
p2.font.size = Pt(28)
p2.font.color.rgb = HUAWEI_RED
p2.alignment = PP_ALIGN.CENTER
p3 = tf.add_paragraph()
p3.text = "输出：答案"
p3.font.size = Pt(18)
p3.font.bold = True
p3.font.color.rgb = BLACK
p3.alignment = PP_ALIGN.CENTER
p4 = tf.add_paragraph()
p4.text = ""
p4.font.size = Pt(10)
p5 = tf.add_paragraph()
p5.text = "程序员手写每一步逻辑\n适合规则明确的问题"
p5.font.size = Pt(14)
p5.font.color.rgb = RGBColor(0x66, 0x66, 0x66)
p5.alignment = PP_ALIGN.CENTER

# ML card
card_ml = slide3.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Inches(6.9), Inches(4.0), Inches(5.8), Inches(3.0)
)
card_ml.fill.solid()
card_ml.fill.fore_color.rgb = WHITE
card_ml.line.color.rgb = HUAWEI_RED
card_ml.line.width = Pt(1)

header_ml = slide3.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Inches(6.9), Inches(4.0), Inches(5.8), Inches(0.5)
)
header_ml.fill.solid()
header_ml.fill.fore_color.rgb = HUAWEI_RED
header_ml.line.fill.background()
txBox = slide3.shapes.add_textbox(Inches(6.9), Inches(4.05), Inches(5.8), Inches(0.4))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "机器学习"
p.font.size = Pt(18)
p.font.bold = True
p.font.color.rgb = WHITE
p.alignment = PP_ALIGN.CENTER

txBox = slide3.shapes.add_textbox(Inches(7.3), Inches(4.7), Inches(5.0), Inches(2.0))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "输入：数据 + 答案"
p.font.size = Pt(20)
p.font.bold = True
p.font.color.rgb = HUAWEI_RED
p.alignment = PP_ALIGN.CENTER
p2 = tf.add_paragraph()
p2.text = "↓"
p2.font.size = Pt(28)
p2.font.color.rgb = HUAWEI_RED
p2.alignment = PP_ALIGN.CENTER
p3 = tf.add_paragraph()
p3.text = "输出：规则（模型）"
p3.font.size = Pt(18)
p3.font.bold = True
p3.font.color.rgb = BLACK
p3.alignment = PP_ALIGN.CENTER
p4 = tf.add_paragraph()
p4.text = ""
p4.font.size = Pt(10)
p5 = tf.add_paragraph()
p5.text = "算法自动发现数据中的模式\n适合规则复杂或未知的问题"
p5.font.size = Pt(14)
p5.font.color.rgb = RGBColor(0x66, 0x66, 0x66)
p5.alignment = PP_ALIGN.CENTER

# ============ PAGE 4: Core Concepts ============
slide4 = prs.slides.add_slide(prs.slide_layouts[6])
add_white_bg(slide4)
add_red_top_line(slide4)
add_title_bar(slide4, "核心思想与对比", "Core Concepts & Comparison", "04")

# Core concept card
card4 = slide4.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(1.6), Inches(12.1), Inches(1.2)
)
card4.fill.solid()
card4.fill.fore_color.rgb = HUAWEI_RED_BG
card4.line.color.rgb = HUAWEI_RED
card4.line.width = Pt(1)

txBox = slide4.shapes.add_textbox(Inches(0.8), Inches(1.7), Inches(11.5), Inches(1.0))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "🎯 核心思想"
p.font.size = Pt(20)
p.font.bold = True
p.font.color.rgb = HUAWEI_RED
p2 = tf.add_paragraph()
p2.text = "让计算机通过经验（数据）自动改善性能，而非每一步都靠人工编写规则。机器学习的本质是从数据中发现规律、提取模式，并利用这些规律对新数据做出预测或决策。"
p2.font.size = Pt(17)
p2.font.color.rgb = BLACK

# Table
txBox = slide4.shapes.add_textbox(Inches(0.6), Inches(3.1), Inches(6), Inches(0.4))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "传统编程与机器学习的全面对比"
p.font.size = Pt(20)
p.font.bold = True
p.font.color.rgb = BLACK

# Table header
table_top = Inches(3.6)
row_h = Inches(0.5)
col_widths = [Inches(2.2), Inches(4.5), Inches(4.5)]
table_left = Inches(1.0)

headers = ["对比维度", "传统编程", "机器学习"]
rows = [
    ["输入方式", "数据 + 规则", "数据 + 答案"],
    ["实现方式", "程序员手写每一步逻辑", "算法自动发现数据中的模式"],
    ["适用场景", "规则明确的问题（如计算器）", "规则复杂或未知的问题（如图像识别）"],
    ["扩展性", "需重写代码适配新场景", "用新数据重新训练即可"],
]

def add_table_cell(slide, left, top, width, height, text, bg_color=None, text_color=BLACK, bold=False, font_size=16, align=PP_ALIGN.LEFT):
    """Add a table cell as a rectangle + text"""
    cell = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, height
    )
    if bg_color:
        cell.fill.solid()
        cell.fill.fore_color.rgb = bg_color
    else:
        cell.fill.background()
    cell.line.color.rgb = HUAWEI_GRAY4
    cell.line.width = Pt(0.5)
    
    txBox = slide.shapes.add_textbox(left + Inches(0.1), top + Inches(0.05), width - Inches(0.2), height - Inches(0.1))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = text_color
    p.alignment = align

# Headers
left_pos = table_left
for i, h in enumerate(headers):
    add_table_cell(slide4, left_pos, table_top, col_widths[i], row_h, h, 
                   bg_color=HUAWEI_RED, text_color=WHITE, bold=True, font_size=16, align=PP_ALIGN.CENTER)
    left_pos += col_widths[i]

# Rows
for row_idx, row in enumerate(rows):
    left_pos = table_left
    y = table_top + row_h * (row_idx + 1)
    for col_idx, cell_text in enumerate(row):
        bg = HUAWEI_GRAY5 if col_idx == 0 else None
        is_red = (row_idx == 3 and col_idx == 2)
        add_table_cell(slide4, left_pos, y, col_widths[col_idx], row_h, cell_text,
                       bg_color=bg, text_color=HUAWEI_RED if is_red else BLACK, 
                       bold=(col_idx == 0 or is_red), font_size=15)
        left_pos += col_widths[col_idx]

# ============ PAGE 5: ML Workflow ============
slide5 = prs.slides.add_slide(prs.slide_layouts[6])
add_white_bg(slide5)
add_red_top_line(slide5)
add_title_bar(slide5, "机器学习基本流程", "ML Workflow", "05")

steps = [
    ("①", "收集数据", "获取原始数据"),
    ("②", "准备/清洗", "去除噪声与缺失值"),
    ("③", "选择模型", "根据任务选取算法"),
    ("④", "训练模型", "用数据拟合模型"),
    ("⑤", "评估模型", "验证效果与精度"),
    ("⑥", "部署使用", "上线到生产环境"),
]

step_w = Inches(1.6)
step_gap = Inches(0.4)
total_w = 6 * step_w + 5 * step_gap + Inches(1.2)
start_x = (SLIDE_W - total_w) / 2 + Inches(0.5)
step_y = Inches(2.2)

for i, (num, title, desc) in enumerate(steps):
    x = start_x + i * (step_w + step_gap)
    
    # Circle with number
    circle = slide5.shapes.add_shape(
        MSO_SHAPE.OVAL, x + Inches(0.3), step_y, Inches(0.85), Inches(0.85)
    )
    circle.fill.solid()
    circle.fill.fore_color.rgb = HUAWEI_RED
    circle.line.fill.background()
    
    txBox = slide5.shapes.add_textbox(x + Inches(0.3), step_y + Inches(0.1), Inches(0.85), Inches(0.65))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = num
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    # Title
    txBox = slide5.shapes.add_textbox(x, step_y + Inches(1.0), step_w, Inches(0.4))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = BLACK
    p.alignment = PP_ALIGN.CENTER
    
    # Description
    txBox = slide5.shapes.add_textbox(x, step_y + Inches(1.35), step_w, Inches(0.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = desc
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(0x66, 0x66, 0x66)
    p.alignment = PP_ALIGN.CENTER

# Arrow connectors (using text)
arrow_y = step_y + Inches(0.35)
for i in range(5):
    x = start_x + (i + 1) * step_w + i * step_gap + Inches(0.1)
    txBox = slide5.shapes.add_textbox(x, arrow_y, step_gap, Inches(0.4))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "→"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = HUAWEI_RED
    p.alignment = PP_ALIGN.CENTER

# Iteration note
txBox = slide5.shapes.add_textbox(Inches(1.0), Inches(3.8), Inches(11), Inches(0.4))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "↻ 评估不达标则返回步骤③调参优化，形成迭代闭环"
p.font.size = Pt(18)
p.font.bold = True
p.font.color.rgb = HUAWEI_RED
p.alignment = PP_ALIGN.CENTER

# Bottom cards
cards_data = [
    ("数据准备阶段", "数据是机器学习的燃料。收集到的原始数据通常包含噪声、缺失值和异常值，需要经过清洗、标准化、特征工程等处理，才能输入模型。"),
    ("模型训练阶段", "选择合适的算法（如决策树、神经网络）后，将训练数据输入模型，模型通过反复迭代自动调整内部参数，使预测结果越来越准确。"),
    ("部署与监控阶段", "训练好的模型部署到生产环境后，需持续监控其性能。随着新数据的产生，需要定期重新训练以保持模型效果。"),
]

card_w = Inches(3.7)
card_h = Inches(2.0)
card_gap = Inches(0.3)
cards_start_x = Inches(0.8)
cards_y = Inches(4.5)

for i, (title, content) in enumerate(cards_data):
    x = cards_start_x + i * (card_w + card_gap)
    
    card = slide5.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, x, cards_y, card_w, card_h
    )
    card.fill.solid()
    card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = HUAWEI_GRAY4
    card.line.width = Pt(1)
    
    txBox = slide5.shapes.add_textbox(x + Inches(0.2), cards_y + Inches(0.15), card_w - Inches(0.4), Inches(0.4))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = BLACK
    
    txBox = slide5.shapes.add_textbox(x + Inches(0.2), cards_y + Inches(0.6), card_w - Inches(0.4), card_h - Inches(0.8))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = content
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(0x66, 0x66, 0x66)
    p.line_spacing = Pt(20)


# ============ PAGE 6: Three Learning Types ============
slide6 = prs.slides.add_slide(prs.slide_layouts[6])
add_white_bg(slide6)
add_red_top_line(slide6)
add_title_bar(slide6, "三大学习类型", "Three Learning Paradigms", "06")

learning_types = [
    ("监督学习", "Supervised Learning", HUAWEI_RED, "⭐ 最常用", 
     '有"正确答案"的训练数据（标签），模型学习从输入到输出的映射关系。',
     ["需要人工标注数据", "预测结果有明确评判标准", "精度通常较高"],
     ["猫狗分类", "房价预测", "垃圾邮件过滤"]),
    ("无监督学习", "Unsupervised Learning", HUAWEI_GRAY1, "🔍 无标签",
     "没有标签的数据，模型自己发现隐藏的结构、模式或聚类关系。",
     ["无需人工标注", "探索数据内在结构", "结果需要人工解读"],
     ["用户分群", "异常检测", "推荐系统"]),
    ("强化学习", "Reinforcement Learning", HUAWEI_RED, "🎮 试错",
     "通过试错 + 奖励/惩罚信号学习最优策略，类似人类学习过程。",
     ["无需预先收集数据", "与环境交互中学习", "追求长期累积回报"],
     ["AlphaGo", "自动驾驶", "机器人控制"]),
]

card_w_6 = Inches(3.7)
card_gap_6 = Inches(0.35)
card_h_6 = Inches(4.8)
start_x_6 = Inches(0.7)
card_y_6 = Inches(1.7)

for i, (title, en_title, header_color, badge, desc, features, apps) in enumerate(learning_types):
    x = start_x_6 + i * (card_w_6 + card_gap_6)
    
    # Card border
    card = slide6.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, x, card_y_6, card_w_6, card_h_6
    )
    card.fill.solid()
    card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = HUAWEI_GRAY4
    card.line.width = Pt(1)
    
    # Header
    header = slide6.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, x, card_y_6, card_w_6, Inches(0.5)
    )
    header.fill.solid()
    header.fill.fore_color.rgb = header_color
    header.line.fill.background()
    
    txBox = slide6.shapes.add_textbox(x, card_y_6 + Inches(0.05), card_w_6, Inches(0.4))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    # Badge
    badge_color = HUAWEI_RED if badge.startswith("⭐") else HUAWEI_GRAY1
    txBox = slide6.shapes.add_textbox(x + Inches(0.2), card_y_6 + Inches(0.6), card_w_6 - Inches(0.4), Inches(0.3))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = badge
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = badge_color
    
    # Description
    txBox = slide6.shapes.add_textbox(x + Inches(0.2), card_y_6 + Inches(0.95), card_w_6 - Inches(0.4), Inches(0.6))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = desc
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(0x66, 0x66, 0x66)
    
    # Features
    feat_box = slide6.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, x + Inches(0.2), card_y_6 + Inches(1.7), card_w_6 - Inches(0.4), Inches(1.2)
    )
    feat_box.fill.solid()
    feat_box.fill.fore_color.rgb = HUAWEI_GRAY5
    feat_box.line.fill.background()
    
    txBox = slide6.shapes.add_textbox(x + Inches(0.35), card_y_6 + Inches(1.75), card_w_6 - Inches(0.7), Inches(1.1))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "核心特点"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = BLACK
    for feat in features:
        p2 = tf.add_paragraph()
        p2.text = "• " + feat
        p2.font.size = Pt(13)
        p2.font.color.rgb = RGBColor(0x66, 0x66, 0x66)
        p2.level = 0
    
    # Apps
    txBox = slide6.shapes.add_textbox(x + Inches(0.2), card_y_6 + Inches(3.1), card_w_6 - Inches(0.4), Inches(0.3))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "应用案例"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = BLACK
    
    app_tags_x = x + Inches(0.2)
    app_tags_y = card_y_6 + Inches(3.45)
    for j, app in enumerate(apps):
        tag = slide6.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, app_tags_x + j * Inches(1.1), app_tags_y, Inches(1.0), Inches(0.3)
        )
        tag.fill.solid()
        tag.fill.fore_color.rgb = HUAWEI_GRAY5
        tag.line.color.rgb = HUAWEI_GRAY4
        tag.line.width = Pt(0.5)
        
        txBox = slide6.shapes.add_textbox(app_tags_x + j * Inches(1.1), app_tags_y + Inches(0.02), Inches(1.0), Inches(0.26))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = app
        p.font.size = Pt(12)
        p.font.color.rgb = RGBColor(0x66, 0x66, 0x66)
        p.alignment = PP_ALIGN.CENTER


# ============ PAGE 7: Applications ============
slide7 = prs.slides.add_slide(prs.slide_layouts[6])
add_white_bg(slide7)
add_red_top_line(slide7)
add_title_bar(slide7, "日常应用场景", "Real-World Applications", "07")

txBox = slide7.shapes.add_textbox(Inches(0.6), Inches(1.6), Inches(11), Inches(0.4))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "机器学习已深入日常生活，在以下场景中发挥着关键作用："
p.font.size = Pt(17)
p.font.color.rgb = RGBColor(0x66, 0x66, 0x66)

app_cards = [
    ("抖音/小红书推荐", "根据用户的点赞、停留时间、评论等行为数据，预测并推荐用户下一个可能喜欢的视频内容。", "算法：协同过滤 + 深度学习"),
    ("人脸解锁", "通过学习用户的面部特征点分布，精确区分\u201c是你\u201d和\u201c不是你\u201d，支持活体检测防伪造。", "算法：卷积神经网络（CNN）"),
    ("语音助手（Siri/小爱）", "将用户的语音信号转为文字（ASR），再理解文字背后的意图（NLU），最后生成响应。", "算法：RNN + Transformer"),
    ("垃圾邮件过滤", "从海量邮件中学习\u201c垃圾邮件\u201d的文本模式、发件特征、链接特征，自动拦截并分类。", "算法：朴素贝叶斯 + 随机森林"),
    ("天气/股票预测", "从大量历史数据中识别时间序列规律（周期性、趋势性），预测未来天气变化或股票走势。", "算法：LSTM + 时间序列分析"),
]

card_w_7 = Inches(5.5)
card_h_7 = Inches(1.4)
card_gap_7 = Inches(0.25)
start_x_7 = Inches(0.7)
start_y_7 = Inches(2.2)

for i, (title, desc, algo) in enumerate(app_cards):
    if i < 3:
        x = start_x_7
        y = start_y_7 + i * (card_h_7 + card_gap_7)
        w = card_w_7
    else:
        x = start_x_7 + card_w_7 + Inches(0.7)
        y = start_y_7 + (i - 3) * (card_h_7 + card_gap_7)
        w = card_w_7
    
    card = slide7.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, x, y, w, card_h_7
    )
    card.fill.solid()
    card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = HUAWEI_GRAY4
    card.line.width = Pt(1)
    
    # Left red border
    border = slide7.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, x, y, Pt(4), card_h_7
    )
    border.fill.solid()
    border.fill.fore_color.rgb = HUAWEI_RED
    border.line.fill.background()
    
    txBox = slide7.shapes.add_textbox(x + Inches(0.25), y + Inches(0.1), w - Inches(0.5), Inches(0.35))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = BLACK
    
    txBox = slide7.shapes.add_textbox(x + Inches(0.25), y + Inches(0.45), w - Inches(0.5), Inches(0.55))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = desc
    p.font.size = Pt(13)
    p.font.color.rgb = RGBColor(0x66, 0x66, 0x66)
    
    txBox = slide7.shapes.add_textbox(x + Inches(0.25), y + Inches(0.95), w - Inches(0.5), Inches(0.3))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = algo
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = HUAWEI_RED


# ============ PAGE 8: Summary ============
slide8 = prs.slides.add_slide(prs.slide_layouts[6])
add_white_bg(slide8)
add_red_top_line(slide8)
add_title_bar(slide8, "总结与展望", "Summary & Outlook", "08")

# Center summary box
summary = slide8.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Inches(2.5), Inches(2.0), Inches(8.0), Inches(2.5)
)
summary.fill.solid()
summary.fill.fore_color.rgb = HUAWEI_RED
summary.line.fill.background()

txBox = slide8.shapes.add_textbox(Inches(2.5), Inches(2.1), Inches(8.0), Inches(0.3))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "ONE SENTENCE SUMMARY"
p.font.size = Pt(14)
p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
p.font.color.brightness = 0.7
p.alignment = PP_ALIGN.CENTER

txBox = slide8.shapes.add_textbox(Inches(2.5), Inches(2.5), Inches(8.0), Inches(0.6))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "机器学习 = 数据 + 算法"
p.font.size = Pt(28)
p.font.bold = True
p.font.color.rgb = WHITE
p.alignment = PP_ALIGN.CENTER

# Decorative line
line8 = slide8.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Inches(5.5), Inches(3.2), Inches(1.0), Pt(2)
)
line8.fill.solid()
line8.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
line8.fill.fore_color.brightness = 0.4
line8.line.fill.background()

txBox = slide8.shapes.add_textbox(Inches(2.5), Inches(3.4), Inches(8.0), Inches(0.6))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "让计算机自己学会做决策\n而不需要人类手写每一条规则"
p.font.size = Pt(20)
p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
p.font.color.brightness = 0.9
p.alignment = PP_ALIGN.CENTER

# Key metrics
metrics = [("3", "大学习类型\n监督·无监督·强化"), ("6", "步标准化流程\n从数据到部署"), ("N", "种应用场景\n遍及日常生活")]

for i, (num, desc) in enumerate(metrics):
    x = Inches(2.5) + i * Inches(2.8)
    y = Inches(4.8)
    
    card = slide8.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, x, y, Inches(2.4), Inches(1.2)
    )
    card.fill.solid()
    card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = HUAWEI_GRAY4
    card.line.width = Pt(1)
    
    txBox = slide8.shapes.add_textbox(x, y + Inches(0.1), Inches(2.4), Inches(0.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = num
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = HUAWEI_RED
    p.alignment = PP_ALIGN.CENTER
    
    txBox = slide8.shapes.add_textbox(x + Inches(0.1), y + Inches(0.55), Inches(2.2), Inches(0.6))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = desc
    p.font.size = Pt(13)
    p.font.color.rgb = RGBColor(0x66, 0x66, 0x66)
    p.alignment = PP_ALIGN.CENTER

# Outlook text
txBox = slide8.shapes.add_textbox(Inches(1.5), Inches(6.2), Inches(10), Inches(0.5))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "随着大模型、多模态技术的快速发展，机器学习正在向通用人工智能（AGI）迈进，未来将在更多领域释放价值。"
p.font.size = Pt(14)
p.font.color.rgb = RGBColor(0x66, 0x66, 0x66)
p.alignment = PP_ALIGN.CENTER

# Thank you
txBox = slide8.shapes.add_textbox(Inches(4.0), Inches(6.7), Inches(5), Inches(0.4))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "THANK YOU · 感谢聆听"
p.font.size = Pt(16)
p.font.color.rgb = RGBColor(0xBF, 0xBF, 0xBF)
p.alignment = PP_ALIGN.CENTER


# ==== Save ====
output_path = os.path.join(os.path.dirname(__file__), "什么是机器学习.pptx")
prs.save(output_path)
print(f"✅ PPTX saved to: {output_path}")
print(f"📊 Size: {os.path.getsize(output_path) / 1024:.1f} KB")
